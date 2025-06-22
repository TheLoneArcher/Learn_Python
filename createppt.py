from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
from pathlib import Path
import os

def add_bullet_slide(prs, title, bullet_points, layout_index=1):
    """Helper function to add a slide with a title and bullet points."""
    bullet_slide_layout = prs.slide_layouts[layout_index]
    slide = prs.slides.add_slide(bullet_slide_layout)
    
    # Title
    title_shape = slide.shapes.title
    if title_shape:
        title_shape.text = title
        # Optional: Adjust title font size if needed
        # title_shape.text_frame.paragraphs[0].font.size = Pt(32)

    # Bullets
    if slide.placeholders and len(slide.placeholders) > 1:
        body_shape = slide.placeholders[1] # Assuming placeholder 1 is the body
        tf = body_shape.text_frame
        tf.clear()  # Clear existing text
        tf.word_wrap = True # Ensure text wraps
        
        for i, point_text in enumerate(bullet_points):
            level = 0
            text = point_text
            if isinstance(point_text, tuple): # (text, level)
                text = point_text[0]
                level = point_text[1]
            
            p = tf.add_paragraph()
            p.text = text
            p.level = level
            # Optional: Adjust bullet font size if needed
            # if i == 0: # Adjust first paragraph, others will follow if not explicitly set
            #     p.font.size = Pt(18)
    else:
        print(f"Warning: Slide '{title}' might not have a body placeholder or has an unexpected layout.")
    return slide

def create_os_presentation():
    prs = Presentation()

    # --- Title Slide (Layout 0) ---
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    title_shape = slide.shapes.title
    subtitle_shape = slide.placeholders[1]

    title_shape.text = "Operating Systems - Unit I Overview"
    subtitle_shape.text = "Annamacharya Institute of Technology & Sciences (Autonomous) - Regulation: AK23"

    # --- Content Slides (Layout 1 is typically title and content) ---

    # Page 1: Introduction, OS Definition, Diagram, Types (Partial)
    add_bullet_slide(prs, "Operating Systems Overview - Introduction", [
        "UNIT-I: Operating Systems",
        "Topics Covered:",
        ("Introduction, OS functions, OS operations, Computing environments", 1),
        ("Free and Open-Source Operating Systems", 1),
        ("System Structures: OS Services, User and OS Interface, system calls, Types of System Calls, system programs", 1),
        ("OS Design and Implementation, OS structure, Building and Booting an OS, OS debugging", 1)
    ])

    add_bullet_slide(prs, "Introduction to Operating Systems (OS)", [
        "Definition 1: An OS is an intermediary between users and computer hardware.",
        "Provides an environment for users to execute programs conveniently and efficiently.",
        "Definition 2 (Technical): Software that manages hardware.",
        "Controls allocation of resources and services (memory, processors, devices, information).",
        "Definition 3: A program acting as an interface between user and hardware, controlling program execution.",
        "(Diagram: Users -> Software (System/Application) -> OS -> Hardware (CPU, RAM, I/O))"
    ])

    add_bullet_slide(prs, "Types Of Operating Systems (1/4)", [
        "Batch Operating System:",
        ("Does not interact with the computer directly.", 1),
        ("An operator groups similar jobs with same requirements into batches.", 1),
        "Time-sharing Operating System:",
        ("Allows many users to share computer resources (maximum utilization).", 1),
        "Distributed Operating System:",
        ("Manages a group of different computers, making them appear as a single computer.", 1),
        ("Designed for networked computers.", 1),
        ("Allows multiple users to access shared resources and communicate.", 1),
        ("Examples: Microsoft Windows Server, various Linux server distributions.", 1)
    ])

    # Page 2: Types (Continued), OS Functions (Intro), Memory Management (Intro)
    add_bullet_slide(prs, "Types Of Operating Systems (2/4)", [
        "Network Operating System:",
        ("Runs on a server.", 1),
        ("Provides capability to manage data, users, groups, security, applications, and networking functions.", 1),
        "Real-time Operating System (RTOS):",
        ("Serves real-time systems; very small time interval for processing and responding to inputs.", 1),
        ("Designed for quick, deterministic responses.", 1),
        ("Applications: embedded systems, industrial control, robotics.", 1)
    ])

    add_bullet_slide(prs, "Types Of Operating Systems (3/4)", [
        "Multiprocessing Operating System:",
        ("Boosts performance of multiple CPUs within a single computer system.", 1),
        ("Multiple CPUs are linked; jobs can be divided and executed more quickly.", 1),
        "Single-User Operating Systems:",
        ("Designed to support a single user at a time.", 1),
        ("Examples: Microsoft Windows (personal), Apple macOS.", 1),
        "Multi-User Operating Systems:",
        ("Designed to support multiple users simultaneously.", 1),
        ("Examples: Linux, Unix.", 1)
    ])
    
    add_bullet_slide(prs, "Types Of Operating Systems (4/4)", [
        "Embedded Operating Systems:",
        ("Designed to run on devices with limited resources (smartphones, wearables, appliances).", 1),
        ("Examples: Google's Android, Apple's iOS.", 1),
        "Cluster Operating Systems:",
        ("Designed to run on a group of computers (a cluster) as a single system.", 1),
        ("Used for high-performance computing and high availability/reliability applications.", 1),
        ("Examples: Rocks Cluster Distribution, OpenMPI.", 1)
    ])

    add_bullet_slide(prs, "Operating Systems Functions - Overview", [
        "Key functions of an OS include:",
        "Memory Management",
        "Processor Management",
        "Device Management",
        "File Management",
        "Security",
        "Control over system performance",
        "Job accounting",
        "Error detecting aids",
        "Coordination between other software and users"
    ])

    add_bullet_slide(prs, "Memory Management", [
        "Manages Primary Memory (Main Memory).",
        "Main memory: large array of words/bytes, each with its own address.",
        "Provides fast storage directly accessible by CPU.",
        "Programs must be in main memory to be executed.",
        "OS activities for memory management:",
        ("Keeps track of primary memory: what part is in use, by whom, what part is not.", 2),
        ("In multiprogramming, OS decides which process gets memory, when, and how much.", 2),
        ("Allocates memory when a process requests it.", 2),
        ("De-allocates memory when a process no longer needs it or is terminated.", 2)
    ])

    # Page 3: Processor, Device, File Management, Other Activities
    add_bullet_slide(prs, "Processor Management", [
        "In a multiprogramming environment, OS decides which process gets the processor, when, and for how much time.",
        "This function is called process scheduling.",
        "OS activities for processor management:",
        ("Keeps track of processor and status of process (task managed by 'traffic controller').", 2),
        ("Allocates the processor (CPU) to a process.", 2),
        ("De-allocates processor when no longer required.", 2)
    ])

    add_bullet_slide(prs, "Device Management", [
        "OS manages device communication via their respective drivers.",
        "OS activities for device management:",
        ("Keeps tracks of all devices (task managed by 'I/O controller').", 2),
        ("Decides which process gets the device, when, and for how much time.", 2),
        ("Allocates the device efficiently.", 2),
        ("De-allocates devices.", 2)
    ])

    add_bullet_slide(prs, "File Management", [
        "File system: normally organized into directories for easy navigation and usage.",
        "Directories may contain files and other directories.",
        "OS activities for file management:",
        ("Keeps track of information, location, uses, status, etc. (collectively known as file system).", 2),
        ("Decides who gets the resources.", 2),
        ("Allocates the resources.", 2),
        ("De-allocates the resources.", 2)
    ])

    add_bullet_slide(prs, "Other Important OS Activities", [
        "Security:",
        ("Preventing unauthorized access to programs and data (e.g., passwords).", 1),
        "Control over system performance:",
        ("Recording delays between service requests and system responses.", 1),
        "Job accounting:",
        ("Keeping track of time and resources used by various jobs and users.", 1),
        "Error detecting aids:",
        ("Production of dumps, traces, error messages, and other debugging/error detecting aids.", 1)
    ])

    # Page 4: OS Objectives/Goals, Operations, Dual-Mode
    add_bullet_slide(prs, "Operating Systems Objectives/Goals", [
        "Coordination between other software and users:",
        ("Assignment of compilers, interpreters, assemblers, etc., to users.", 1),
        "Convenience (Primary Goal):",
        ("Simplify user interaction. System programs (utilities) assist in program creation, file management, I/O control.", 1),
        "Efficiency (Secondary Goal):",
        ("Efficient operation of the system. Manages resources (data movement, storage, processing).", 1),
        ("Kernel/nucleus in main memory contains frequently used functions.", 1),
        ("Determines processor time allocation for efficient resource utilization.", 1),
        "Ability to Evolve:",
        ("Constructed to permit effective development, testing, and introduction of new functions.", 1),
        ("Reasons for evolution: Hardware upgrades, new services, bug fixes.", 1)
    ])

    add_bullet_slide(prs, "Operating-System Operations", [
        "Events are signaled by interrupts or traps.",
        "Trap (or exception): software-generated interrupt.",
        ("Caused by an error (e.g., division by zero, invalid memory access).", 1),
        ("Or by a specific request from a user program for an OS service (system call).", 1),
        "For each interrupt type, specific OS code (interrupt service routine) determines action.",
    ])

    add_bullet_slide(prs, "Dual-Mode and Multimode Operation", [
        "Need for two separate modes: user mode and kernel mode.",
        "Kernel mode also called: supervisor mode, system mode, or privileged mode.",
        "Mode bit in hardware indicates current mode: kernel (0) or user (1).",
        "User application executes in user mode.",
        "When requesting OS service (via system call), system transitions from user to kernel mode.",
        "(Diagram: Transition from user to kernel mode via system call and trap)"
    ])

    # Page 5: Dual-Mode (Cont.), Timer
    add_bullet_slide(prs, "Dual-Mode Operation (Continued)", [
        "System boots in kernel mode; OS loads and starts user applications in user mode.",
        "Hardware switches to kernel mode on trap or interrupt.",
        "OS always gains control in kernel mode.",
        "Purpose: Protects OS from errant users and users from each other.",
        "Achieved by designating harmful instructions as privileged instructions.",
        "Privileged instructions execute only in kernel mode.",
        ("Attempt to execute in user mode results in a trap to OS.", 1),
        "Modes can extend beyond two (e.g., virtualization - VMM mode).",
        ("VMM has more privileges than user, fewer than kernel.", 1),
        "Lack of hardware-supported dual mode (e.g., MS-DOS on Intel 8088) causes shortcomings:",
        ("User program can wipe out OS.", 1),
        ("Multiple programs can write to a device simultaneously.", 1),
        "Modern Intel CPUs provide dual-mode operation."
    ])

    add_bullet_slide(prs, "Timer", [
        "Prevents user programs from getting stuck (infinite loop) or not returning control to OS.",
        "Timer can be set to interrupt the computer after a specified period.",
        "Period can be fixed (e.g., 1/60 second) or variable (e.g., 1ms to 1s).",
        "Variable timer: often implemented with a fixed-rate clock and a counter.",
        ("OS sets counter; clock ticks decrement counter; interrupt occurs when counter reaches 0.", 1),
        "Example: 10-bit counter with 1ms clock allows interrupts from 1ms to 1024ms.",
        "Usage: Initialize counter with allowed program run time (e.g., 7 min = 420s counter). Timer interrupts, counter decremented."
    ])

    # Page 6: Computing Environments
    add_bullet_slide(prs, "Computing Environments - Traditional", [
        "Stand-alone general-purpose machines (single desktop computer for a single user, local applications).",
        "Without LAN or WAN connection initially.",
        "Most systems now connect to Internet via web portals (share printers, scanners).",
        "Networking is ubiquitous; home systems use firewalls for protection.",
        "(Image: Desktop Computer)"
    ])

    add_bullet_slide(prs, "Computing Environments - Mobile", [
        "Mobile computers interconnect via wireless networks.",
        "Includes handheld smartphones, tablets, etc.",
        "Extra OS features (e.g., GPS).",
        "Allows new app types (augmented reality, virtual reality).",
        "Uses IEEE 802.11 wireless protocols or cellular data networks.",
        "Leaders: Apple iOS, Google Android."
    ])
    
    add_bullet_slide(prs, "Computing Environments - Client-Server", [
        "Distributed application structure.",
        "Partitions tasks/workloads between servers (providers of resource/service) and clients (service requesters).",
        "Compute-server system:",
        ("Provides interface for client to request services (e.g., database).", 1),
        "File-server system:",
        ("Provides interface for clients to store and retrieve files.", 1),
        "(Image: Client-Server interaction diagram)"
    ])

    # Page 7: Distributed Systems, Client-Server Computing (Diagrams)
    add_bullet_slide(prs, "Computing Environments - Distributed Systems", [
        "Consists of multiple software components on multiple computers, running as a single system.",
        "Computers can be physically close (connected by LAN) or geographically distant (connected by WAN).",
        "(Diagram: LAN and WAN examples)"
    ])
    
    add_bullet_slide(prs, "Client-Server Computing (Revisited)", [
        "Modern systems often act as server systems satisfying requests from client systems.",
        "Specialized distributed system structure.",
        "General structure:",
        ("(Diagram: Server connected to client desktop, laptop, smartphone via Network)",1),
        "Server systems broadly categorized as:",
        ("Compute servers",1),
        ("File servers",1)
    ])


    # Page 8: Peer-to-Peer, Cloud Computing
    add_bullet_slide(prs, "Computing Environments - Peer-to-Peer (P2P)", [
        "Another model of distributed system.",
        "Does not distinguish clients and servers; all nodes are peers.",
        "Each node may act as client, server, or both.",
        "Joining P2P network:",
        ("Register service with central lookup service, OR", 1),
        ("Broadcast request for service and respond via discovery protocol.", 1),
        "Popularized in late 1990s (Napster, Gnutella for file sharing).",
        "Example: Skype (VoIP - voice over IP).",
        "(Diagram: Interconnected P2P clients)"
    ])

    add_bullet_slide(prs, "Computing Environments - Cloud Computing", [
        "Cloud service provider delivers services to clients across a network.",
        "Services: computing, storage, even apps as a service.",
        "Utilizes virtualization technology to provide services.",
        "(Diagram: Cloud architecture with user, enterprise, provider)"
    ])

    # Page 9: Cloud Computing (Types, Services), Real-Time Embedded
    add_bullet_slide(prs, "Cloud Computing - Types & Services", [
        "Types of Clouds:",
        ("Public Cloud: Available via Internet to anyone willing to pay.", 1),
        ("Private Cloud: Run by a company for its own use.", 1),
        ("Hybrid Cloud: Includes both public and private components.", 1),
        "Cloud Services Models:",
        ("Software as a Service (SaaS): Applications via Internet (e.g., word processor).", 1),
        ("Platform as a Service (PaaS): Software stack for application use via Internet (e.g., database server).", 1),
        ("Infrastructure as a Service (IaaS): Servers or storage via Internet (e.g., backup storage).", 1),
        "(Diagram: Application/Platform/Infrastructure layers in cloud)"
    ])

    add_bullet_slide(prs, "Real-Time Embedded Systems", [
        "Embedded system responding to real-time situations with embedded software/hardware within specified time constraints.",
        "Different from general-purpose computers.",
        "Vary considerably: special purpose, limited functionalities, real-time OS (RTOS).",
        "Some have OSes, some perform tasks without an OS.",
        "RTOS has well-defined fixed time constraints.",
        "Processing must be done within constraint.",
        "(Diagram: Hard Real Time vs Soft Real Time examples - Aircraft, DVD Player)"
    ])

    # Page 10: Open-Source OS, Linux
    add_bullet_slide(prs, "Open-Source Operating Systems", [
        "Available in source-code format (not just compiled binary).",
        "Examples: Linux (famous open-source), MS Windows (closed-source).",
        "Hybrid approach: Apple's Mac OS X / iOS (open-source Darwin kernel, proprietary components).",
        "Source code allows modification and recompilation (learning tool).",
        "Reverse engineering binaries is difficult; comments are lost.",
        "Benefits:",
        ("Community debugging, analysis, support, suggestions.", 1),
        ("Potentially more secure (more eyes on code).", 1),
        ("Bugs found and fixed faster.", 1)
    ])

    add_bullet_slide(prs, "Open-Source OS - Linux Example", [
        "GNU project: produced UNIX-compatible tools (compilers, editors) but no kernel.",
        "1991: Linus Torvalds (Finland) released rudimentary UNIX-like kernel using GNU tools.",
        "Invited worldwide contributions via Internet.",
        "Rapid growth due to weekly updates and thousands of programmers.",
        "Resulting GNU/Linux OS has spawned hundreds of unique distributions."
    ])

    # Page 11: Linux Distributions, BSD UNIX, Solaris
    add_bullet_slide(prs, "Linux Distributions & Running Linux", [
        "Distributions (custom builds): RedHat, SUSE, Fedora, Debian, Slackware, Ubuntu.",
        "Vary in: function, utility, installed apps, hardware support, UI, purpose.",
        ("Example: RedHat Enterprise Linux (commercial), PCLinuxOS (LiveCD/LiveDVD).", 1),
        "Running Linux on Windows (example approach):",
        ("1. Download free VMware Player.", 1),
        ("2. Choose a Linux 'appliance' (virtual machine image) from VMware.", 1),
        ("3. Boot the virtual machine in VMware Player.", 1)
    ])

    add_bullet_slide(prs, "Open-Source OS - BSD UNIX", [
        "Started 1978 as derivative of AT&T's UNIX.",
        "UC Berkeley (UCB) releases required AT&T license (not truly open source initially).",
        "Development slowed by AT&T lawsuit.",
        "1994: Fully functional, open-source 4.4BSD-lite released.",
        "Distributions: FreeBSD, NetBSD, OpenBSD, DragonflyBSD.",
        "Source code (e.g., FreeBSD) available in /usr/src/ (kernel in /usr/src/sys).",
        "Darwin (Mac OS X core kernel) is based on BSD UNIX, open-sourced by Apple."
    ])
    
    add_bullet_slide(prs, "Open-Source OS - Solaris", [
        "Commercial UNIX-based OS by Sun Microsystems.",
        "Originally SunOS (BSD based), moved to AT&T System V UNIX base in 1991.",
        "2005: Sun open-sourced most Solaris code as OpenSolaris project.",
        "Oracle's purchase of Sun (2009) left project state unclear.",
        "Original 2005 source still available.",
        "Project Illumos: forked from OpenSolaris base, expanded features, basis for several products."
    ])


    # Page 12: Operating System Services (List & Diagram)
    add_bullet_slide(prs, "Operating System Services - Overview", [
        "OS services provided for programmer convenience, to make programming easier.",
        "An OS provides an environment for program execution.",
        "Services cater to:",
        ("Programs (System programs)", 1),
        ("Users", 1),
        "Common classes of services (vary by OS):",
        ("User Interface", 1),
        ("Program Execution", 1),
        ("I/O Operations", 1),
        ("File Systems Manipulation", 1),
        ("Communication", 1),
        ("Error Management/Detection", 1),
        ("Resource Allocation", 1),
        ("Accounting", 1),
        ("Protection and Security", 1),
        "(Diagram: OS services architecture - User/Programs -> UI -> System Calls -> Services -> OS -> Hardware)"
    ])

    # Page 13: User Interface, Program Execution, I/O Operations, File System Manipulation
    add_bullet_slide(prs, "OS Services - User Interface (UI)", [
        "Almost all OS have a UI.",
        "Types of UI:",
        ("Command-Line Interface (CLI): Text commands via keyboard (e.g., Bash, CMD).", 1),
        ("Batch Interface: Commands & directives in files, then executed.", 1),
        ("Graphical User Interface (GUI): Window system, pointing device (mouse), menus, keyboard (e.g., Windows Desktop, macOS Finder).", 1),
        "Some systems offer multiple UI variations.",
        "(Diagram: UI types - GUI, batch, command line)"
    ])

    add_bullet_slide(prs, "OS Services - Program Execution", [
        "OS loads the program into memory to execute.",
        "CPU executes the program (now a process).",
        "Process termination:",
        ("Success message on normal completion.", 1),
        ("Failure message on error.", 1),
        "(Diagram showing program execution block among other services)"
    ])

    add_bullet_slide(prs, "OS Services - I/O Operations", [
        "Running programs may require I/O (file or I/O device).",
        "Special functions for specific devices:",
        ("Saving file to disk.", 1),
        ("Providing input data.", 1),
        ("Blanking a display screen.", 1),
        "Users usually cannot control I/O devices directly (for efficiency/protection).",
        "Users interact with I/O devices via OS (device drivers)."
    ])

    add_bullet_slide(prs, "OS Services - File System Manipulation", [
        "Programs need to read and write files and directories.",
        "Need to create, delete, search for files, and list file information.",
        "Some OS include permission management (allow/deny access based on ownership)."
    ])

    # Page 14: Communication, Error Management, Resource Allocation, User & OS Interface Intro
    add_bullet_slide(prs, "OS Services - Communication", [
        "Processes interact to interchange information.",
        "Processes can run on the same computer or different computers on a network.",
        "Implementation methods:",
        ("Shared Memory: Two or more processes read/write to a shared memory section.", 1),
        ("Message Passing: Packets of information (predefined formats) moved between processes by OS.", 1)
    ])

    add_bullet_slide(prs, "OS Services - Error Management (Detection)", [
        "OS continuously monitors system for errors.",
        "Errors can occur in:",
        ("CPU and memory hardware.", 1),
        ("I/O devices.", 1),
        ("User programs.", 1),
        "OS should take appropriate action for each error type for correct and consistent computing.",
        "Actions: Halt system, terminate error-causing process, or return error code."
    ])

    add_bullet_slide(prs, "OS Services - Resource Allocation", [
        "When multiple users or jobs run concurrently, resources must be allocated.",
        "OS manages various resource types.",
        "Special allocation code for: CPU cycles, main memory, file storage.",
        "General request/release code for: I/O devices.",
        "CPU scheduling routines consider CPU speed, jobs, registers, etc.",
        "Routines also allocate printers, USB drives, other peripherals."
    ])

    add_bullet_slide(prs, "User and Operating-System Interface", [
        "Serves as a bridge between user and the underlying OS.",
        "Allows users and applications to:",
        ("Communicate with the OS.", 1),
        ("Request services.", 1),
        ("Interact with hardware/software resources effectively.", 1),
        "Categorized based on how users/applications interact with OS."
    ])
    
    add_bullet_slide(prs, "Types of User Interfaces (1/2)", [
        "1. Command-Line Interface (CLI):",
        ("Description: Users interact by typing text-based commands in a terminal/shell.", 2),
        ("Examples: Bash (Linux), Command Prompt (Windows), PowerShell.", 2),
        ("Advantages: Lightweight, efficient, allows scripting for automation, direct OS functionality access.", 2),
        ("Disadvantages: Steeper learning curve for beginners, limited visual feedback.", 2),
        "2. Graphical User Interface (GUI):",
        ("Description: Users interact via graphical elements (windows, icons, buttons, menus).", 2),
        ("Examples: Windows Desktop, macOS Finder, GNOME/KDE on Linux.", 2),
        ("Advantages: User-friendly, intuitive, rich visual feedback, easier navigation, suitable for non-technical users.", 2),
        ("Disadvantages: Resource-intensive, limited for advanced configurations compared to CLI.", 2)
    ])

    # Page 15: User Interfaces (Cont.), System Calls Intro
    add_bullet_slide(prs, "Types of User Interfaces (2/2)", [
        "3. Touchscreen Interfaces:",
        ("Description: Common on mobile devices/tablets; interaction via touch gestures (swiping, pinching, tapping).", 2),
        ("Examples: Android, iOS, Windows with touch support.", 2),
        ("Advantages: Highly intuitive, portable, optimized for small screens.", 2),
        ("Disadvantages: Not suitable for detailed configurations/coding, input can be less precise.", 2),
        "4. Voice User Interfaces (VUI):",
        ("Description: Users interact via voice commands.", 2),
        ("Examples: Siri (iOS), Cortana (Windows), Google Assistant.", 2),
        ("Advantages: Hands-free interaction, ideal for accessibility and smart devices.", 2),
        ("Disadvantages: Limited accuracy with complex commands, dependent on voice recognition quality.", 2)
    ])

    add_bullet_slide(prs, "System Calls", [
        "How a program requests a service from an OS's kernel.",
        "Provides an interface to OS services.",
        "A means for user/application programs to call upon OS services.",
        "Generally written in C or C++, some in assembly for performance.",
        "Example: Sequence of system calls to copy a file (Acquire input/output file names, open input, create output, loop read/write, close files, terminate)."
    ])

    # Page 16: System Calls (API, Parameter Passing)
    add_bullet_slide(prs, "System Calls - API and Invocation", [
        "Most programmers use Application Programming Interface (API) instead of direct low-level system calls.",
        "API benefits: Greater program portability between systems.",
        "API makes appropriate system calls via the system call interface (often using a table lookup for numbered calls).",
        "(Diagram: User app -> open() -> System Call Interface -> Kernel's open() implementation -> return)",
        "Parameter Passing to System Calls:",
        ("Generally via registers.", 1),
        ("Less commonly, values pushed onto the stack.", 1),
        ("Large data blocks: accessed indirectly via memory address (passed in register or on stack).", 1),
        "(Diagram: Passing parameters as a table via register X)"
    ])

    # Page 17: Types of System Calls
    add_bullet_slide(prs, "Types of System Calls (Categories)", [
        "Grouped into six major categories:",
        "1. Process Control:",
        ("end, abort; load, execute; create/terminate process; get/set process attributes; wait for time/event; signal event; allocate/free memory.", 1),
        "2. File Manipulation (Management):",
        ("create/delete file; open, close; read, write, reposition; get/set file attributes.", 1),
        "3. Device Manipulation (Management):",
        ("request/release device; read, write, reposition; get/set device attributes; logically attach/detach devices.", 1),
        "4. Information Maintenance:",
        ("get/set time or date; get/set system data; get/set process, file, or device attributes.", 1),
        "5. Communications:",
        ("create/delete communication connection; send/receive messages; transfer status info; attach/detach remote devices.", 1),
        "6. Protection:",
        ("(Controls access to resources - not explicitly detailed in diagram but implied by security functions)",1)
    ])
    
    # Page 18: Examples of System Calls (Win/Unix), Standard C Library
    add_bullet_slide(prs, "Examples of Windows and Unix System Calls", [
        "Process Control:",
        ("Windows: CreateProcess(), ExitProcess(), WaitForSingleObject()", 1),
        ("Unix: fork(), exit(), wait()", 1),
        "File Manipulation:",
        ("Windows: CreateFile(), ReadFile(), WriteFile(), CloseHandle()", 1),
        ("Unix: open(), read(), write(), close()", 1),
        "Device Manipulation:",
        ("Windows: SetConsoleMode(), ReadConsole(), WriteConsole()", 1),
        ("Unix: ioctl(), read(), write()", 1),
        "Information Maintenance:",
        ("Windows: GetCurrentProcessID(), SetTimer(), Sleep()", 1),
        ("Unix: getpid(), alarm(), sleep()", 1),
        "Communication:",
        ("Windows: CreatePipe(), CreateFileMapping(), MapViewOfFile()", 1),
        ("Unix: pipe(), shmget(), mmap()", 1),
        "Protection:",
        ("Windows: SetFileSecurity(), InitializeSecurityDescriptor()", 1),
        ("Unix: chmod(), umask(), chown()", 1)
    ])

    add_bullet_slide(prs, "Example of Standard C Library Interaction", [
        "Standard C library provides part of the system-call interface for many UNIX/Linux versions.",
        "Example: C program invokes `printf()`.",
        "C library intercepts `printf()` call.",
        "Library invokes necessary system call(s) in OS (e.g., `write()`).",
        "Library takes value returned by `write()` and passes it back to user program.",
        "(Diagram: User program -> printf() -> C Library -> write() syscall -> kernel's write() -> return path)"
    ])
    
    # Page 19: Process Control Details (DOS vs UNIX)
    add_bullet_slide(prs, "System Calls - Process Control Details", [
        "Includes: end, abort, load, execute, create/terminate process, get/set attributes, wait, signal, allocate/free memory.",
        "Processes: created, launched, monitored, paused, resumed, stopped.",
        "If one process pauses/stops, another may be launched/resumed.",
        "Abnormal stop: may need core dumps or diagnostics.",
        "Comparison: DOS (single-tasking) vs. UNIX (multi-tasking).",
        "DOS Execution:",
        ("Command interpreter unloads itself partially to free memory.", 1),
        ("Loads process, transfers control. Interpreter resumes after process completion.", 1),
        ("(Diagram: MS-DOS execution - (a) startup, (b) running program)", 1),
        "UNIX Execution:",
        ("Command interpreter remains resident.", 1),
        ("User can switch to interpreter, place process in background.", 1),
        ("Interpreter uses `fork()`: creates a child process (clone of interpreter - parent).", 1),
        ("Child process uses `exec()`: replaces its code with the desired program.", 1),
        ("Parent (interpreter) can wait for child or issue new prompt immediately (child runs in background).", 1)
    ])

    # Page 20: File & Device Management, Info Maintenance, Communication Details
    add_bullet_slide(prs, "System Calls - File Management Details", [
        "Includes: create/delete file, open, close, read, write, reposition, get/set attributes.",
        "Supports directories and ordinary files.",
        "Directory structure can be implemented using ordinary files or other means."
    ])

    add_bullet_slide(prs, "System Calls - Device Management Details", [
        "Includes: request/release device, read, write, reposition, get/set attributes, attach/detach.",
        "Devices can be physical (disk drives) or virtual/abstract (files, partitions, RAM disks).",
        "Some systems represent devices as special files (e.g., /dev on UNIX); accessing 'file' calls device drivers."
    ])

    add_bullet_slide(prs, "System Calls - Information Maintenance Details", [
        "Includes: get/set time, date, system data, process/file/device attributes.",
        "May provide: memory dump, single-step program execution, program tracing (for debugging)."
    ])

    add_bullet_slide(prs, "System Calls - Communication Details", [
        "Includes: create/delete connection, send/receive messages, transfer status, attach/detach remote devices.",
        "Message Passing Model must support:",
        ("Identify remote process/host.", 1),
        ("Establish connection.", 1),
        ("Open/close connection.", 1),
        ("Transmit messages.", 1),
        ("Wait for incoming messages (blocking/non-blocking).", 1),
        ("Delete connection.", 1)
    ])
    
    # Page 21: Shared Memory, Protection, System Programs
    add_bullet_slide(prs, "Communication - Shared Memory Model", [
        "Shared Memory Model must support:",
        ("Create and access memory shared among processes (and threads).", 1),
        ("Provide locking mechanisms for simultaneous access.", 1),
        ("Free up shared memory / dynamically allocate it.", 1),
        "Comparison:",
        ("Message passing: simpler, easier (esp. inter-computer), good for small data.", 1),
        ("Shared memory: faster, better for large data (esp. read-mostly or few writers).", 1)
    ])

    add_bullet_slide(prs, "System Calls - Protection", [
        "Provides mechanisms to control user/process access to system resources.",
        "System calls adjust access mechanisms.",
        "Allow non-privileged users temporary elevated permissions under controlled circumstances.",
        "Crucial for all systems today (due to network connectivity), not just multi-user systems."
    ])

    add_bullet_slide(prs, "System Programs (Utilities / Applications)", [
        "Provide OS functionality via separate applications (not part of kernel/command interpreter).",
        "Most systems ship with useful applications (calculators, editors like Notepad).",
        "Categories:",
        ("File management: create, delete, copy, rename, print, list files/directories.", 1),
        ("Status information: date, time, users, running processes, data logging, system registries (config).", 1),
        ("File modification: text editors, tools to change file contents.", 1),
        ("Programming-language support: compilers, linkers, debuggers, profilers, assemblers, library management, interpreters, make.", 1),
        ("Program loading and execution: loaders (static, dynamic, overlay), interactive debuggers.", 1),
        ("Communications: mail, web browsers, remote logins, file transfers, remote command execution.", 1),
        ("Background services (daemons): started at boot, run continuously (network daemons, print servers, schedulers, error monitoring).", 1),
        "Many OS also include application programs for additional services (copying files, checking time/date)."
    ])

    # Page 22: OS Design & Implementation, Goals, Mechanisms & Policies
    add_bullet_slide(prs, "Operating-System Design and Implementation", [
        "Outline: 1. Design Goals, 2. Mechanisms and Policies, 3. Implementation.",
        "OS allows user applications to interact with system hardware.",
        "OS itself doesn't provide functions but an environment for useful work.",
        "Internal structure of OS varies widely."
    ])

    add_bullet_slide(prs, "Design Goals - Requirements", [
        "Affected by choice of hardware and type of system (batch, time-sharing, single/multi-user, distributed, real-time, general purpose).",
        "Requirements divided into two basic groups:",
        ("User goals: Convenient to use, easy to learn, reliable, safe, fast.", 1),
        ("System goals: Easy to design, create, implement, maintain; flexible, reliable, error-free, efficient.", 1),
        "Specifying and designing an OS is a highly creative software engineering task."
    ])
    
    add_bullet_slide(prs, "Requirements of Designing Operating System", [
        "No unique solution to defining OS requirements.",
        "Wide range of systems -> different requirements -> variety of solutions -> different environments.",
        "Examples: Batch, time-sharing, single/multi-user, distributed, real-time, general purpose.",
        "Collect requirements based on user demand, system type, and environment."
    ])

    add_bullet_slide(prs, "Mechanisms and Policies", [
        "Policy: What needs to be done? (Conceptual, e.g., 'Interrupt after every 100 seconds').",
        "Mechanism: How to do something? (Implementation, e.g., 'timer hardware').",
        "Important principle: Separation of policy from mechanism.",
        "Policies may change over time; a general mechanism should require few changes if policy changes."
    ])

    # Page 23: Implementation Languages, OS Structure Intro
    add_bullet_slide(prs, "OS Implementation", [
        "Once designed, OS must be implemented.",
        "OS are collections of many programs, written by many people over time.",
        "Difficult to make general statements about implementation specifics."
    ])
    
    add_bullet_slide(prs, "Languages used to implement OS", [
        "Early OS: Assembly language (e.g., MS-DOS for 8088 processor).",
        "Modern OS: Higher-level languages like C, C++.",
        "Lowest kernel levels might still use assembly.",
        "Can be written in multiple languages:",
        ("Higher-level routines in C.", 1),
        ("System programs in C/C++, scripting (Perl, Python), shell scripts.", 1),
        "Emulators: Duplicate functionality of one system on another, allow OS on non-native hardware.",
        "Linux OS programs: Written in all these languages."
    ])

    add_bullet_slide(prs, "Advantages & Disadvantages of Higher-Level Language for OS", [
        "Advantages:",
        ("Code written faster.", 1),
        ("More compact code.", 1),
        ("Easier to understand and debug.", 1),
        ("Easier to port (move to other hardware).", 1),
        "Disadvantages:",
        ("Reduced speed.", 1),
        ("Increased storage requirements.", 1)
    ])

    add_bullet_slide(prs, "Operating-System Structure - Overview", [
        "Four different structures based on how components are interconnected into a kernel:",
        "1. Simple Structure",
        "2. Layered Approach",
        "3. Microkernels",
        "4. Modules"
    ])

    add_bullet_slide(prs, "OS Structure - Simple Structure (MS-DOS)", [
        "Many OS do not have well-defined structures.",
        "Typically small, simple, limited systems.",
        "Example: MS-DOS.",
        "Interfaces and levels of functionality not well separated.",
        "Application programs can access basic I/O routines directly (display, disk drives).",
        "Vulnerability: Malicious programs can cause entire system crashes.",
        "MS-DOS also limited by hardware (Intel 8088: no dual mode, no hardware protection).",
        "(Diagram: MS-DOS layer structure - App -> Resident System Prog -> MS-DOS drivers -> ROM BIOS)"
    ])

    # Page 24: UNIX Structure, Layered Approach
    add_bullet_slide(prs, "OS Structure - Original UNIX (Monolithic)", [
        "Used a simple structure (less layered than ideal).",
        "Almost all OS components were in one large layer (kernel).",
        "Not well broken down into layered subsystems.",
        "This monolithic structure was difficult to implement and maintain.",
        "(Diagram: Traditional UNIX system structure - Users -> Shells/Compilers/Libraries -> Syscall Interface -> Kernel components -> Hardware Interface -> Hardware)"
    ])

    add_bullet_slide(prs, "OS Structure - Layered Approach", [
        "Breaks OS into smaller layers.",
        "Each layer rests on the layer below it.",
        "Relies solely on services from the next lower layer.",
        "Allows independent development and debugging of layers (assuming lower layers are trusted).",
        "Problem: Deciding layer order (no layer can call higher layers; chicken-and-egg situations).",
        "Efficiency: Can be less efficient; requests filter through multiple layers to reach hardware.",
        "(Diagram: A layered operating system - concentric circles from Hardware (Layer 0) to User Interface (Layer N))"
    ])

    # Page 25: Microkernels
    add_bullet_slide(prs, "OS Structure - Microkernels", [
        "Basic idea: Remove all non-essential services from the kernel.",
        "Implement non-essential services as system applications.",
        "Goal: Make kernel as small and efficient as possible.",
        "Most microkernels provide:",
        ("Basic process and memory management.", 1),
        ("Message passing between other services.", 1),
        "Benefits:",
        ("Enhanced security/protection (most services in user mode).", 1),
        ("Easier system expansion (add system apps, not rebuild kernel).", 1),
        "Examples:",
        ("Mach: First widely known, now part of Mac OS X.", 1),
        ("Windows NT (original): Was microkernel, suffered performance issues. NT 4.0+ more monolithic.", 1),
        ("QNX: Real-time OS for embedded systems.", 1),
        "(Diagram: Architecture of a typical microkernel - User mode components communicating via microkernel)"
    ])

    # Page 26: Modules, Hybrid Systems
    add_bullet_slide(prs, "OS Structure - Modules", [
        "Modern OS development is object-oriented.",
        "Relatively small core kernel + set of dynamically linkable modules.",
        "Example: Solaris (Diagram: Solaris loadable modules - scheduling, file systems, etc. around core kernel).",
        "Similarity to layers: Each subsystem has clearly defined tasks/interfaces.",
        "Difference: Any module can contact any other (eliminates multi-layer traversal problem).",
        "Kernel is relatively small (like microkernels) but doesn't need to implement message passing (modules contact directly)."
    ])

    add_bullet_slide(prs, "OS Structure - Hybrid Systems", [
        "Combine different structures to address performance, security, usability.",
        "Example: Linux and Solaris.",
        ("Monolithic for performance (single address space).", 1),
        ("Modular for dynamic addition of new functionality.", 1),
        "Example: Windows.",
        ("Largely monolithic (for performance).", 1),
        ("Retains some microkernel-like behavior (e.g., separate subsystems/personalities as user-mode processes).", 1),
        ("Supports dynamically loadable kernel modules.", 1)
    ])

    # Page 27: Mac OS X, iOS Structure
    add_bullet_slide(prs, "Hybrid Systems - Mac OS X", [
        "Uses a hybrid structure; layered system.",
        "Top Layers: Aqua user interface, application environments/services (Java, Cocoa, QuickTime, BSD).",
        "(Diagram: The Mac OS X structure)",
        "Cocoa environment: API for Objective-C (Mac OS X applications).",
        "Kernel Environment (below UI/App layers):",
        ("Mach microkernel.", 1),
        ("BSD UNIX kernel components.", 1),
        "Mach provides: memory management, RPCs, IPC (message passing), thread scheduling.",
        "BSD component provides: command-line interface, networking, file systems, POSIX APIs (Pthreads).",
        "I/O Kit: For device driver development, dynamically loadable modules (kernel extensions)."
    ])

    add_bullet_slide(prs, "Hybrid Systems - iOS", [
        "Mobile OS by Apple for iPhone, iPad.",
        "Structured on Mac OS X, with added mobile-specific functionality.",
        "Does not directly run Mac OS X applications.",
        "(Diagram: Architecture of Apple's iOS - Core OS, Core Services, Media Services, Cocoa Touch)",
        "Cocoa Touch: Objective-C API for iOS apps.",
        ("Differs from Cocoa: supports hardware features unique to mobile (e.g., touch screens).", 1),
        "Media Services: Graphics, audio, video services.",
        "Core Services: Cloud computing support, databases, etc.",
        "Core Operating System (bottom layer): Based on Mac OS X kernel environment."
    ])

    # Page 28: Android Structure
    add_bullet_slide(prs, "Hybrid Systems - Android", [
        "Designed by Open Handset Alliance (led by Google) for smartphones/tablets.",
        "Runs on various mobile platforms; open-sourced.",
        "Layered software stack for mobile app development.",
        "(Diagram: Architecture of Google's Android - Linux Kernel, Libraries/Android Runtime, App Framework, Apps)",
        "Bottom Layer: Linux kernel (modified by Google, currently outside normal Linux releases).",
        "Android Runtime Environment:",
        ("Core set of libraries.", 1),
        ("Dalvik Virtual Machine (DVM).", 1),
        "App Development: Java language, but uses separate Android API (not standard Java API).",
        "Execution: Java class files -> Java byte code -> translated to executable for DVM.",
        "DVM: Designed for Android, optimized for limited memory/CPU of mobile devices.",
        "Libraries: Frameworks for web browsers (webkit), database (SQLite), multimedia.",
        "libc library: Similar to standard C library but smaller (for slower mobile CPUs)."
    ])

    # Page 29: Operating-System Debugging, Failure Analysis, Performance Tuning
    add_bullet_slide(prs, "Operating-System Debugging", [
        "Activity of finding and fixing errors (bugs) in a system (hardware and software)."
    ])

    add_bullet_slide(prs, "Debugging - Failure Analysis", [
        "Process failure: OS writes error info to a log file.",
        "OS can take a core dump (capture of process memory) for later analysis.",
        "Debuggers probe running programs and core dumps to explore code/memory.",
        "Kernel debugging: More complex than user-level.",
        ("Due to kernel size, complexity, hardware control, lack of user-level tools.", 1),
        "Kernel failure: Called a crash.",
        ("Error info saved to log file; memory state saved to crash dump.", 1),
        "OS debugging and process debugging often use different tools/techniques."
    ])

    add_bullet_slide(prs, "Debugging - Performance Tuning", [
        "Seeks to improve performance by removing processing bottlenecks.",
        "Requires monitoring system performance.",
        "OS methods for performance data:",
        ("Trace listings: Logs interesting events with time/parameters to a file. Analysis program processes log.", 1),
        ("Interactive tools: Query state of system components for bottlenecks.", 2),
        ("Example: UNIX `top` command (displays resource usage, top processes).", 2),
        ("Other tools: Display disk I/O, memory allocation, network traffic.", 2),
        "Windows Task Manager: Similar tool for Windows (current apps, processes, CPU/memory usage, network stats).",
        "(Diagram: Screenshot of Windows Task Manager)"
    ])

    # Page 30: DTrace
    add_bullet_slide(prs, "Debugging - DTrace", [
        "Facility that dynamically adds probes to a running system (user processes and kernel).",
        "Probes queried via D programming language to get info about kernel, system state, process activities.",
        "Example: Tracing `ioctl()` system call, shows function calls within kernel.",
        ("Lines ending 'U': user mode; 'K': kernel mode.", 1),
        "(Diagram: Solaris 10 dtrace output following a system call)",
        "Essential for debugging user-level and kernel code interactions."
    ])

    # Page 31: DTrace (Continued)
    add_bullet_slide(prs, "DTrace - Functionality", [
        "Kernel debugging was historically difficult before DTrace (Solaris 10).",
        ("CPU breakpoint feature halts system, unsuitable for multiuser OS kernels.", 1),
        "DTrace runs on production systems with minimal impact; resets system to pre-debugging state after use.",
        "Components:",
        ("Compiler, Framework, Providers (of probes), Consumers (of probes).", 1),
        "Providers create probes; kernel tracks probes (hash table).",
        "Enabling a probe: Rewrites code to call `dtrace_probe()`, then original operation.",
        "DTrace compiler generates 'safe' byte code run in kernel.",
        "Privileges: Only 'root' or DTrace-privileged users can use it (can access/modify private kernel data).",
        "Consumer requests probes; probe fires -> emits data managed by kernel.",
        "Enabling Control Blocks (ECBs): Kernel actions when probes fire. Predicates can filter ECBs.",
        "Termination: Consumer terminates -> ECBs removed. No consumers -> probe removed (original code restored).",
        "Example DTrace code: Record CPU time for processes with UID 101."
    ])

    # Page 32: DTrace Output, System Boot
    add_bullet_slide(prs, "DTrace - Example Output & Availability", [
        "(Diagram: Output of DTrace script showing processes and CPU time in nanoseconds)",
        "DTrace is part of open-source OpenSolaris (Solaris 10).",
        "Added to other OS (Mac OS X, FreeBSD) with compatible licenses.",
        "Spreading due to unique capabilities."
    ])

    add_bullet_slide(prs, "System Boot (Building and Booting an OS)", [
        "After OS generation, it must be made available to hardware.",
        "Booting: Procedure of starting a computer by loading the kernel.",
        "Bootstrap program (or bootstrap loader):",
        ("Small piece of code, locates kernel, loads it into main memory, starts execution.", 1),
        ("Usually in Read-Only Memory (ROM) - needs no initialization, virus-resistant.", 1),
        "Bootstrap program tasks:",
        ("Run diagnostics (check machine state).", 1),
        ("If diagnostics pass, continue booting.", 1),
        ("Initialize system aspects (CPU registers, device controllers, main memory).", 1),
        ("Starts the OS.", 1)
    ])

    # Page 33: System Boot (Firmware, Large OS Boot)
    add_bullet_slide(prs, "System Boot - Firmware and Variations", [
        "Small systems (phones, tablets, game consoles): Entire OS in ROM.",
        ("Suitable for small OS, simple hardware, rugged operation.", 1),
        ("Problem: Changing bootstrap code requires changing ROM chips.", 1),
        "EPROM (Erasable Programmable Read-Only Memory): Read-only, but writable with specific command.",
        "Firmware: All forms of ROM (characteristics between hardware/software).",
        ("Problem: Slower execution than RAM. Some systems copy OS from firmware to RAM for speed.", 1),
        ("Problem: Relatively expensive, so usually small amounts available.", 1)
    ])

    add_bullet_slide(prs, "System Boot - Large OS (Disk-based)", [
        "For large OS (Windows, Mac OS X, UNIX) or frequently changing systems:",
        ("Bootstrap loader in firmware; OS on disk.", 1),
        ("Bootstrap runs diagnostics, reads a fixed block (e.g., block zero - boot block) from disk into memory, executes it.", 1),
        ("Boot block program may be sophisticated enough to load entire OS and start execution.", 1),
        "GRUB: Example open-source bootstrap program for Linux systems.",
        "Disk-bound bootstrap & OS easily changed by writing new versions to disk.",
        "Boot disk (or system disk): Disk that has a bootable OS.",
        "Full bootstrap program (loaded from boot block) traverses file system, finds kernel, loads to memory, starts execution.",
        "System is considered 'running' at this point."
    ])

    # --- Save Presentation ---
    try:
        # Get the path to the user's Downloads folder
        downloads_path = Path.home() / "Downloads"
        # Ensure the Downloads directory exists (it usually does, but good practice)
        os.makedirs(downloads_path, exist_ok=True) 
        
        file_name = "Operating_Systems_Unit1_Summary.pptx"
        pptx_path = downloads_path / file_name
        
        prs.save(pptx_path)
        print(f"Presentation successfully created and saved to: {pptx_path}")
        return True
    except Exception as e:
        print(f"Error creating or saving presentation: {e}")
        # Fallback: Save to current directory if Downloads folder fails
        try:
            fallback_path = Path(file_name)
            prs.save(fallback_path)
            print(f"Presentation successfully created and saved to fallback location: {fallback_path.resolve()}")
            return True
        except Exception as fe:
            print(f"Error saving presentation to fallback location: {fe}")
            return False

if __name__ == '__main__':
    create_os_presentation()

